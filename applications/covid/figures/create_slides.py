#!/usr/bin/env python
"""Create PowerPoint slides summarizing COVID scRNA-seq datasets."""

import anndata
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import io

DATA_DIR = os.path.join(os.path.dirname(__file__), "..", "data")
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "dataset_summary.pptx")

DATASETS = {
    "combat": {
        "title": "COMBAT Consortium",
        "citation": "COMBAT Consortium, Cell (2022)",
        "tissue": "Whole blood (PBMCs)",
    },
    "ren": {
        "title": "Ren et al.",
        "citation": "Ren et al., Cell (2021)",
        "tissue": "PBMCs + airway samples",
    },
    "stevenson": {
        "title": "Stephenson et al.",
        "citation": "Stephenson et al., Nature Medicine (2021)",
        "tissue": "Whole blood (PBMCs)",
    },
    "combined": {
        "title": "Combined (All Datasets)",
        "citation": "COMBAT + Ren + Stephenson",
        "tissue": "Whole blood (PBMCs) + airway samples",
    },
}


def make_histogram(adata, title):
    """Create a cells-per-donor histogram with log x-scale, truncated at 50k."""
    counts = adata.obs["donor_id"].value_counts().values
    # Clip at 50000 for display
    counts_clipped = np.clip(counts, None, 50000)

    fig, ax = plt.subplots(figsize=(6, 3.5))
    bins = np.logspace(np.log10(counts_clipped.min()), np.log10(50000), 21)
    ax.hist(counts_clipped, bins=bins, edgecolor="white", color="#4C72B0")
    ax.set_xscale("log")
    ax.set_xlabel("Cells per donor", fontsize=12)
    ax.set_ylabel("Number of donors", fontsize=12)
    ax.set_title(title, fontsize=13)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


def add_summary_slide(prs, name, info, adata):
    """Add a slide summarizing one dataset."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = info["title"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Citation
    p2 = tf.add_paragraph()
    p2.text = info["citation"]
    p2.font.size = Pt(14)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # Stats table
    assays = ", ".join(adata.obs["assay"].unique().tolist())
    n_cells = f"{adata.n_obs:,}"
    n_donors = adata.obs["donor_id"].nunique()
    donors_per_label = adata.obs.groupby("label", observed=True)["donor_id"].nunique()

    lines = [
        ("Tissue", info["tissue"]),
        ("scRNA-seq technology", assays),
        ("Total cells (post-QC)", n_cells),
        ("Total donors", str(n_donors)),
        ("  Control", str(donors_per_label.get("control", 0))),
        ("  Mild", str(donors_per_label.get("mild", 0))),
        ("  Severe", str(donors_per_label.get("severe", 0))),
    ]

    table_top = Inches(1.6)
    table = slide.shapes.add_table(
        len(lines), 2, Inches(0.5), table_top, Inches(4.5), Inches(0.35 * len(lines))
    ).table
    table.columns[0].width = Inches(2.3)
    table.columns[1].width = Inches(2.2)

    for i, (label, value) in enumerate(lines):
        for j, text in enumerate([label, value]):
            cell = table.cell(i, j)
            cell.text = text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                if label.startswith("  "):
                    paragraph.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # Histogram
    hist_buf = make_histogram(adata, f"Cells per donor — {info['title']}")
    slide.shapes.add_picture(hist_buf, Inches(5.2), Inches(1.6), Inches(4.5))


DL_BENCHMARKS = [
    {
        "title": "CellCnn",
        "citation": "Arvaniti & Claassen, Nature Communications (2017)",
        "bullets": [
            "Detects rare disease-associated cell subsets via representation learning",
            "1D convolutions (kernel_size=1) applied independently to each cell",
            "Top-k mean pooling selects the most informative cells per filter",
            "Pooled representation fed to a linear classifier",
            "Designed for cases where only a small fraction of cells carry the signal",
        ],
    },
    {
        "title": "scAGG",
        "citation": "Verlaan et al., CSBJ (2025)",
        "bullets": [
            "Single-cell gene expression analysis using graph-based aggregation",
            "2-layer MLP cell encoder with ELU activations",
            "Mean pooling aggregates cell embeddings into a sample-level representation",
            "NoGraph variant (MLP + mean pooling) performs on par with graph variants",
            "Simple architecture serves as a strong MIL-style baseline",
        ],
    },
    {
        "title": "ScRAT",
        "citation": "Mao et al., Bioinformatics (2024)",
        "bullets": [
            "Transformer-based method originally developed for single-cell multi-omic integration",
            "Linear projection + sinusoidal positional encoding → Transformer encoder",
            "Mean pooling over cell tokens followed by a classification head",
            "Cell-type-aware sample mixup augmentation for regularization",
            "Uses BCE loss (one-vs-rest) with cosine LR schedule and warmup",
        ],
    },
]


def add_benchmark_slide(prs, info):
    """Add a slide describing one DL benchmark method."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Benchmark: {info['title']}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Citation
    p2 = tf.add_paragraph()
    p2.text = info["citation"]
    p2.font.size = Pt(14)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # Bullet points
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5))
    btf = body.text_frame
    btf.word_wrap = True
    for i, bullet in enumerate(info["bullets"]):
        if i == 0:
            bp = btf.paragraphs[0]
        else:
            bp = btf.add_paragraph()
        bp.text = bullet
        bp.font.size = Pt(18)
        bp.space_after = Pt(10)
        bp.level = 0
        # Bullet character
        bp.text = f"\u2022  {bullet}"


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for name, info in DATASETS.items():
        path = os.path.join(DATA_DIR, f"{name}_processed.h5ad")
        print(f"Loading {path}...")
        adata = anndata.read_h5ad(path)
        add_summary_slide(prs, name, info, adata)

    for info in DL_BENCHMARKS:
        add_benchmark_slide(prs, info)

    prs.save(OUTPUT_PATH)
    print(f"Saved to {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
